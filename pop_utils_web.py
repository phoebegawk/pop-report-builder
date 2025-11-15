from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io
from typing import List, Tuple

# Paths
SCRIPT_DIR = Path(__file__).resolve().parent if "__file__" in globals() else Path(".").resolve()
ASSETS_DIR = SCRIPT_DIR / "assets"

TEMPLATE_PATH = ASSETS_DIR / "PoP's Report - Template (Date).pptx"
LOGO_PATH = ASSETS_DIR / "GAWK LOGO (PURPLE).png"
BACKGROUND_PATH = ASSETS_DIR / "Background.jpg"

# Colours
PURPLE = RGBColor(0x54, 0x2D, 0x54)
GAWK_GREEN = RGBColor(0xD7, 0xDF, 0x23)


def convert_date(date_str: str) -> str:
    """Convert DDMMYY to 'MONTH YYYY' in uppercase."""
    try:
        return datetime.strptime(date_str, "%d%m%y").strftime("%B %Y").upper()
    except Exception:
        return "INVALID DATE"


def display_date(date_str: str) -> str:
    """Convert DDMMYY to DD/MM/YY for on-slide display."""
    try:
        return datetime.strptime(date_str, "%d%m%y").strftime("%d/%m/%y")
    except Exception:
        return "INVALID DATE"


def parse_filename(path_like) -> dict | None:
    """
    Parse filenames in the format:
    'Site Name - Site Code - Client - Campaign - DDMMYY - Type[ - OptionalSuffix]'.
    Only the first 6 parts are meaningful for the PoP logic.
    """
    # Accept either a Path or an object with a .name attribute (e.g. UploadedFile)
    if hasattr(path_like, "name"):
        name = Path(path_like.name)
    else:
        name = Path(path_like)

    parts = name.stem.split(" - ")
    if len(parts) < 6:
        # Too short to be valid
        return None

    try:
        site_name = f"{parts[0].strip()} - {parts[1].strip()}"
        client = parts[2].strip()
        campaign = parts[3].strip()
        live_date_raw = parts[4].strip()
        return {
            "site_name": site_name,
            "client": client,
            "campaign": campaign,
            "live_date": live_date_raw,
            "month_year": convert_date(live_date_raw),
            "live_date_display": display_date(live_date_raw),
        }
    except Exception:
        return None


def extract_live_date_priority(path: Path):
    """
    Sort key matching desktop script:
    1. By actual live date (earliest first)
    2. Then by base filename (first 5 parts)
    3. Then by suffix priority: Cam (0) < Mock (1) < Others (2)
    """
    try:
        name = path.name.lower()
        parts = path.stem.split(" - ")
        base_key = " - ".join(parts[:5]).lower()
        suffix = name.split(" - ")[-1]

        if "cam" in suffix:
            suffix_priority = 0
        elif "mock" in suffix:
            suffix_priority = 1
        else:
            suffix_priority = 2

        date_str = parts[4]
        date = datetime.strptime(date_str, "%d%m%y")
        return (date, base_key, suffix_priority)
    except Exception:
        return (datetime.min, path.name.lower(), 9)


def _add_front_slide_content(prs: Presentation, first_info: dict) -> None:
    first_slide = prs.slides[0]
    for shape in first_slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text
        if "Client Name" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = first_info["client"]
            r.font.name = "Montserrat"
            r.font.size = Pt(60)
            r.font.bold = True
            r.font.color.rgb = RGBColor(255, 255, 255)
        elif "DATE" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = first_info["month_year"]
            r.font.name = "Montserrat"
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = RGBColor(255, 255, 255)
        elif "Campaign Name" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = first_info["campaign"]
            r.font.name = "Montserrat"
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = GAWK_GREEN


def _add_text(slide, x, y, w, h, text, size=23, bold=True, color=RGBColor(255, 255, 255)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Montserrat"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def generate_presentation_bytes(image_paths: List[Path]) -> Tuple[bytes, str]:
    """
    Core PoP logic reproduced from the desktop GUI script,
    but saving to an in-memory bytes object instead of using file dialogs.
    """
    # Filter & normalise
    image_files: list[Path] = []
    for p in image_paths:
        p = Path(p)
        if p.is_dir():
            image_files.extend(
                f for f in p.glob("*") if f.suffix.lower() in [".jpg", ".jpeg", ".png"]
            )
        elif p.suffix.lower() in [".jpg", ".jpeg", ".png"]:
            image_files.append(p)

    image_files = sorted(image_files, key=extract_live_date_priority)

    if not image_files:
        raise FileNotFoundError("No JPG, JPEG, or PNG files found for PoP generation.")

    first_info = parse_filename(image_files[0])
    if not first_info:
        raise ValueError("First image filename is invalid. Cannot determine client/campaign/date.")

    prs = Presentation(TEMPLATE_PATH)
    blank_layout = prs.slide_layouts[5]

    _add_front_slide_content(prs, first_info)

    for img_path in image_files:
        details = parse_filename(img_path)
        if not details:
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Background strip
        slide.shapes.add_picture(
            str(BACKGROUND_PATH),
            Cm(0),
            Cm(-0.01),
            width=Cm(29.7),
            height=Cm(21),
        )

        # Main PoP image, scaled to fit the 24.89 x 12.87cm box at (3.4, 4.45)
        img = Image.open(img_path)
        iw, ih = img.size
        img_aspect = iw / ih
        box_aspect = 24.89 / 12.87

        if img_aspect > box_aspect:
            new_width = Cm(24.89)
            new_height = Cm(24.89 / img_aspect)
        else:
            new_height = Cm(12.87)
            new_width = Cm(12.87 * img_aspect)

        slide.shapes.add_picture(
            str(img_path),
            Cm(3.4),
            Cm(4.45),
            width=new_width,
            height=new_height,
        )

        # Gawk logo top-right
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Cm(23.8),
            Cm(1.52),
            width=Cm(4.49),
            height=Cm(1.46),
        )

        # Gawk green vertical strip
        rect = slide.shapes.add_shape(
            1,  # MSO_SHAPE_RECTANGLE
            Cm(0),
            Cm(0),
            Cm(1.22),
            Cm(21),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = GAWK_GREEN
        rect.line.fill.background()

        # Vertical 'PROOF OF POSTING'
        tb = slide.shapes.add_textbox(Cm(-9.73), Cm(9.96), Cm(20.71), Cm(0.94))
        tb.rotation = 270
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "PROOF OF POSTING"
        r.font.name = "Montserrat"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = PURPLE

        # Site + Live Date labels/values
        _add_text(
            slide,
            Cm(3.06),
            Cm(2.5),
            Cm(3.09),
            Cm(1.24),
            "Site:",
            color=GAWK_GREEN,
        )
        _add_text(
            slide,
            Cm(5.36),
            Cm(2.5),
            Cm(13.25),
            Cm(1.24),
            details["site_name"],
        )
        _add_text(
            slide,
            Cm(3.06),
            Cm(18),
            Cm(4.76),
            Cm(1.24),
            "Live Date:",
            color=GAWK_GREEN,
        )
        _add_text(
            slide,
            Cm(7.79),
            Cm(18),
            Cm(12.35),
            Cm(1.24),
            details["live_date_display"],
        )

    # Remove example slide (index 1) and move 'Gotta love rectangles' to the end
    if len(prs.slides) >= 3:
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[1])
        final_slide = prs.slides._sldIdLst[1]
        prs.slides._sldIdLst.remove(final_slide)
        prs.slides._sldIdLst.append(final_slide)

    safe_client = first_info["client"].strip()
    safe_month = first_info["month_year"]
    output_name = f"PoP Report - {safe_client} ({safe_month}).pptx"

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read(), output_name


def generate_presentation_from_uploads(uploaded_files) -> Tuple[bytes, str]:
    """
    Wrapper for Streamlit: takes a list of UploadedFile objects and returns
    (pptx_bytes, suggested_filename).
    """
    temp_dir = SCRIPT_DIR / "uploaded_pop_images"
    temp_dir.mkdir(exist_ok=True)

    paths: list[Path] = []
    for uf in uploaded_files:
        dest = temp_dir / uf.name
        with open(dest, "wb") as f:
            f.write(uf.getbuffer())
        paths.append(dest)

    return generate_presentation_bytes(paths)